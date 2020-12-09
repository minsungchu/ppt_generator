import sys
from PyQt5.QtWidgets import QApplication, QWidget, QPushButton, QToolTip, QGridLayout, QLabel, QTextEdit
from PyQt5.QtGui import QIcon
from PyQt5.QtGui import QFont
from pptx import Presentation  # 라이브러리
from pptx.util import Inches  # 사진, 표등을 그리기 위해


class MyApp(QWidget):

    def __init__(self):
        super().__init__()
        self.btnConvert = QPushButton("입력 내용을 PPT로 변환(&S)")
        self.texteditLyrics = QTextEdit()
        self.labelLyrics = QLabel("찬양 가사")
        self.texteditTitle = QTextEdit()
        self.labelTitle = QLabel("찬양 제목")
        self.initUI()

    def initUI(self):
        """툴팁 설정"""
        QToolTip.setFont(QFont('SansSerif', 10))

        """위젯 아이템"""

        """그리드 레이아웃 설정"""
        grid = QGridLayout()
        self.setLayout(grid)
        grid.addWidget(self.labelTitle, 0, 0)
        grid.addWidget(self.texteditTitle, 1, 0)
        grid.addWidget(self.labelLyrics, 2, 0)
        grid.addWidget(self.texteditLyrics, 3, 0)
        grid.addWidget(self.btnConvert, 4, 0)

        """Stylesheet"""
        self.setStyleSheet(
            "background-color : #303745;"
            "color : #CFD6C9;"
        )

        self.texteditTitle.setStyleSheet(
            """
            QTextEdit{
                min-height : 50px;
                max-height : 50px;
            }
            """
        )

        self.btnConvert.setStyleSheet(
            """
            QPushButton{
                min-height : 50px;
                max-height : 50px;
                font-size  : 20px;
            }
            """
        )

        """다이얼로그 제목,아이콘,위치,크기 설정"""
        self.setWindowTitle('스트리밍용 가사 PPT 생성 프로그램')
        self.setWindowIcon(QIcon('schana_icon.jpg'))
        self.move(200, 100)
        self.resize(1000, 800)
        self.show()

        """이벤트 핸들러 등록"""
        self.btnConvert.clicked.connect(self.btnConvertClicked)

    def btnConvertClicked(self):
        pptFile = Presentation("./template.pptx")

        ccm_title = self.texteditTitle.toPlainText()
        ccm_lyrics_raw = self.texteditLyrics.toPlainText().split("\n")
        lyrics_length = len(ccm_lyrics_raw)
        ccm_lyrics = []

        for idx in range(0, lyrics_length):
            if idx == lyrics_length - 1:
                if idx % 2 == 0:
                    ccm_lyrics.append(ccm_lyrics_raw[idx])
                    break
            if idx % 2 == 0:
                if ccm_lyrics_raw[idx + 1] == "":
                    ccm_lyrics.append(ccm_lyrics_raw[idx])
                else:
                    ccm_lyrics.append(ccm_lyrics_raw[idx] + "\n" + ccm_lyrics_raw[idx + 1])

        for idx in range(0, len(ccm_lyrics)):
            slide = pptFile.slides[idx]
            if slide.shapes[0].name == "title":
                title = slide.shapes[0]
                title.text = ccm_title
                lyrics = slide.shapes[1]
                lyrics.text = ccm_lyrics[idx]
            else:
                title = slide.shapes[1]
                title.text = ccm_title
                lyrics = slide.shapes[0]
                lyrics.text = ccm_lyrics[idx]

        pptFile.save("../" + ccm_title + ".pptx")


if __name__ == '__main__':
    app = QApplication(sys.argv)
    ex = MyApp()
    sys.exit(app.exec_())
